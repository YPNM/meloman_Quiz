from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import db_init
from random import shuffle


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

def sorting_table(game_id):
    teamsInGame = db_init.TeamsInGame()
    scoresModel = db_init.ScoresDB()
    roundsModel = db_init.RoundsDB()
    scoresDictionary = {}

    teams = teamsInGame.get_all_teams_in_game(game_id, active=True)
    allscores = scoresModel.get_scores_by_game_id(game_id=game_id)
    rounds = roundsModel.get_all_rounds(game_id=game_id)

    # template
    for round in rounds:
        for score in allscores:
            if (score[4] not in scoresDictionary.keys()):
                scoresDictionary[score[4]] = {'total': 0, 'items': []}
            if (round[0] == score[0]):
                scoresDictionary[score[4]]['total'] += score[3]
                scoresDictionary[score[4]]['items'].append(score)

    maxRounds = 0
    # find empty rounds
    for key, value in scoresDictionary.items():
        if (len(value['items']) < len(rounds)):
            existRounds = []
            team_id = ''
            for item in value['items']:
                existRounds.append(item[0])
                team_id = item[2]
            for round in rounds:
                if (round[0] not in existRounds):
                    value['items'].append([round[0], round[1], team_id, 0, key])
            maxRounds = len(existRounds)
    if(maxRounds == 0):
        maxRounds = len(rounds)


    # sorting by ascending order
    sortedDict = dict(sorted(scoresDictionary.items(), key=lambda x: x[1]['total'], reverse=True))

    if (maxRounds > 3):
        submitDict = {}
        temp = {
            'max': 0,
            'key': '',
        }
        for key, value in sortedDict.items():
            if (key not in submitDict.keys()):
                if (value['total'] == temp['max']):
                    length = len(value['items']) - 1
                    secondLength = len(sortedDict[key]['items']) - 1
                    while length != -1:
                        if value['items'][length][3] > scoresDictionary[temp['key']]['items'][secondLength][3]:
                            res = dict()
                            for secondKey in submitDict.keys():
                                if (secondKey == temp['key']):
                                    res[key] = {
                                        'total': value['total'],
                                        'items': value['items']
                                    }

                                res[secondKey] = submitDict[secondKey]
                            submitDict = res
                            length = -1
                        elif (value['items'][length][3] == scoresDictionary[temp['key']]['items'][secondLength][3]):
                            if (length == 0 and secondLength == 0):
                                submitDict[key] = {
                                    'total': value['total'],
                                    'items': value['items']
                                }
                                length = -1
                            else:
                                length -= 1
                                secondLength -= 1
                            continue
                        elif (value['items'][length][3] < scoresDictionary[temp['key']]['items'][secondLength][3]):
                            submitDict[key] = {
                                'total': value['total'],
                                'items': value['items']
                            }
                            length = -1
                else:
                    temp['max'] = value['total']
                    temp['key'] = key
                    submitDict[key] = {
                        'total': value['total'],
                        'items': value['items']
                    }
    else:
        submitDict = sortedDict
    # add new teams without scores
    for team in teams:
        if (team[3] not in submitDict.keys()):
            submitDict[team[3]] = {
                'total': None,
                'id': team[1]
            }
    return submitDict

def create_powerpoint_with_table(folder, game_id):
    # Создаем новую презентацию

    submitDict = sorting_table(game_id)

    prs = Presentation('table_preset.pptx')
    game_result = {}

    for key, value in submitDict.items():
        game_result[key] = []
        for i in range(0, len(value['items'])):
            arr = value['items'][i]
            arr.append(value['total'])
            game_result[key].append(arr)

    rows = len(game_result.keys()) + 1
    cols = len(game_result.get(next(iter(game_result.keys())))) + 3
    table_length = Inches(4) + Inches(1) + (cols - 2) * Inches(0.7)

    commands = list(game_result.keys())
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            try:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if 'место' in run.text:
                            continue
                        else:
                            if i < 3:
                                try:
                                    run.text = str(commands[2 - i]).upper()
                                except IndexError:
                                    pass
                            else:
                                break
            except AttributeError:
                pass

    # Заканчиваю сортировку -------------------------------------------------------------------------------------------

    # Добавляем первый слайд с общей слайд с таблицей
    slide_layout = prs.slide_layouts[5]  # 5 - это слайд с черным фоном
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Черный цвет фона
    if rows <= 15:
        font_size = Pt(22)
    elif rows <= 18:
        font_size = Pt(18)
    else:
        font_size = Pt(14)
    if 7.5 / rows < 0.5:
        row_height = 7.5 / rows
    else:
        row_height = 0.5

    table = slide.shapes.add_table(rows=rows, cols=cols, left=Inches(13.3 / 2) - table_length / 2,
                                   top=Inches((7.5 / 2) - (rows * row_height / 2)),
                                   width=Inches(9),
                                   height=Inches(row_height) * rows)

    for cell in iter_cells(table.table):
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

    counter = 1

    for i in range(rows):
        for j in range(cols):
            cell = table.table.cell(i, j)
            if i == 0:
                if j == 0:
                    cell.text = "№"
                elif j == 1:
                    cell.text = "Название команды"
                elif j == cols - 1:
                    cell.text = "Всего"
                else:
                    cell.text = str(counter)
                    counter += 1
            else:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)  # Белый цвет текста
                if j == 0:
                    cell.text = str(i)
                elif j == 1:
                    cell.text = str(commands[i - 1])
                elif j == cols - 1:
                    cell.text = str(float(game_result.get(commands[i - 1])[0][5]))
                else:
                    cell.text = str(game_result.get(commands[i - 1])[j - 2][3])
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.size = font_size
            cell.text_frame.paragraphs[0].font.bold = True

    for i in range(cols):
        if i == 1:
            table.table.columns[i].width = Inches(4)
        elif i == cols - 1:
            table.table.columns[i].width = Inches(1)
        else:
            table.table.columns[i].width = Inches(0.7)

    # Создаю второй слайд с промежуточной таблицей ---------------------------------------------------------------------

    shuffled_list = list(range(len(game_result.keys())))
    shuffle(shuffled_list)
    slide_layout = prs.slide_layouts[5]  # 5 - это слайд с черным фоном
    slide_2 = prs.slides.add_slide(slide_layout)
    background = slide_2.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Черный цвет фона

    if len(game_result.get(commands[0])) >= 7:

        table_2_length = Inches(4) + 4 * Inches(0.7)

        table_2 = slide_2.shapes.add_table(rows=rows, cols=5, left=Inches(13.3 / 2) - table_2_length / 2,
                                           top=Inches((7.5 / 2) - (rows * row_height / 2)),
                                           width=Inches(9),
                                           height=Inches(row_height) * rows)

        counter = 5
        for i in range(rows):
            for j in range(5):
                cell = table_2.table.cell(i, j)
                if i == 0:
                    if j == 0:
                        cell.text = "№"
                    elif j == 1:
                        cell.text = "Название команды"
                    else:
                        cell.text = str(counter)
                        counter += 1
                else:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255)  # Белый цвет текста
                    if j == 0:
                        cell.text = str(i)
                    elif j == 1:
                        cell.text = str(commands[shuffled_list[i - 1]])
                    else:
                        cell.text = str(game_result.get(commands[shuffled_list[i - 1]])[j + 4 - 2][3])
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.size = font_size
                cell.text_frame.paragraphs[0].font.bold = True

        for cell in iter_cells(table_2.table):
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

        for i in range(5):
            if i == 1:
                table_2.table.columns[i].width = Inches(4)
            else:
                table_2.table.columns[i].width = Inches(0.7)

    path = f'{folder}/table.pptx'
    # Сохраняем презентацию
    prs.save(path)
    return path

# create_powerpoint_with_table('E:/Programming/WEB/meloman_Quiz', 'b8bfea64-4a3e-11ee-8a2c-6c2408aface6')
